from pptx import Presentation


functions = {
    "random.sample":
        """random.sample = sample(population, k) method of random.Random instance
    Chooses k unique random elements from a population sequence or set.
    
    Returns a new list containing elements from the population while
    leaving the original population unchanged.  The resulting list is
    in selection order so that all sub-slices will also be valid random
    samples.  This allows raffle winners (the sample) to be partitioned
    into grand prize and second place winners (the subslices).
    
    Members of the population need not be hashable or unique.  If the
    population contains repeats, then each occurrence is a possible
    selection in the sample.
    
    To choose a sample in a range of integers, use range as an argument.
    This is especially fast and space efficient for sampling from a
    large population:   sample(range(10000000), 60)""",
    "random.shuffle":
        """random.shuffle = shuffle(x, random=None) method of random.Random instance
    Shuffle list x in place, and return None.
    
    Optional argument random is a 0-argument function returning a
    random float in [0.0, 1.0); if it is the default None, the
    standard random.random will be used.""",
    "random.randint":
        """random.randint = randint(a, b) method of random.Random instance
    Return random integer in range [a, b], including both end points.""",
    "random.gauss":
        """random.gauss = gauss(mu, sigma) method of random.Random instance
    Gaussian distribution.
    
    mu is the mean, and sigma is the standard deviation.  This is
    slightly faster than the normalvariate() function.
    
    Not thread-safe without a lock around calls.""",
    "random.gammavariate":
        """random.gammavariate = gammavariate(alpha, beta) method of random.Random instance
    Gamma distribution.  Not the gamma function!
    
    Conditions on the parameters are alpha > 0 and beta > 0.
    
    The probability distribution function is:
    
                x ** (alpha - 1) * math.exp(-x / beta)
      pdf(x) =  --------------------------------------
                  math.gamma(alpha) * beta ** alpha"""}

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]

for key, val in functions.items():
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = key
    subtitle.text = val

prs.save('test.pptx')
